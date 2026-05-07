#!/usr/bin/env python3
"""
📚 EduSave Bot — O'quv materiallarini kategoriyalarda saqlash boti
  + 🗜 Fayl siqish (rasm va hujjatlar)
  + 📦 Kategoriyani ZIP arxiv qilish
"""

import logging
import sqlite3
import os
import io
import json
import zipfile
import threading
import requests
from datetime import datetime

from PIL import Image
from flask import Flask

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup
from telegram.ext import (
    Application, CommandHandler, MessageHandler,
    CallbackQueryHandler, filters, ContextTypes
)

logging.basicConfig(
    format='%(asctime)s - %(levelname)s - %(message)s',
    level=logging.INFO
)
logger = logging.getLogger(__name__)

DB_PATH        = "edusave.db"
ITEMS_PER_PAGE = 5

# ─────────────────────────── ADMIN ────────────────────────────────
# Admin user ID — faqat siz admin paneli ko'ra olasiz
ADMIN_ID = 7855999182  # ← Sizning Telegram user ID

# ─────────────────────────── GEMINI AI ──────────────────────────
# Gemini API kalit — Render Environment Variables'dan olinadi
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")
GEMINI_URL = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent"

def generate_presentation_content(topic: str) -> dict:
    """Gemini API yordamida prezentatsiya matnini yaratadi."""
    if not GEMINI_API_KEY:
        raise Exception("GEMINI_API_KEY o'rnatilmagan!")

    prompt = f"""Quyidagi mavzu bo'yicha 5 slaydli prezentatsiya yarating: "{topic}"

Faqat JSON formatida javob bering, boshqa hech qanday matn yo'q.
JSON struktura:
{{
  "title_main": "Asosiy sarlavha (1-2 so'z, masalan: Sun'iy Intellekt)",
  "title_sub": "Qo'shimcha sarlavha (1-3 so'z)",
  "intro_question": "Kirish savoli (3-5 so'z, masalan: Bu nima?)",
  "intro_definition": "Asosiy ta'rif (1 jumla, 15-25 so'z)",
  "intro_extra": "Qo'shimcha izoh (1 jumla, 10-15 so'z)",
  "stat1_value": "Birinchi statistika qiymati (masalan: 75%)",
  "stat1_label": "Statistika izohi (3-7 so'z)",
  "stat2_value": "Ikkinchi statistika qiymati (masalan: $1.5T)",
  "stat2_label": "Statistika izohi (3-7 so'z)",
  "directions": [
    {{"icon": "🧠", "title": "1-yo'nalish", "desc": "Tavsif (8-12 so'z)"}},
    {{"icon": "🗣️", "title": "2-yo'nalish", "desc": "Tavsif (8-12 so'z)"}},
    {{"icon": "👁️", "title": "3-yo'nalish", "desc": "Tavsif (8-12 so'z)"}},
    {{"icon": "🤖", "title": "4-yo'nalish", "desc": "Tavsif (8-12 so'z)"}}
  ],
  "applications": [
    {{"icon": "🏥", "title": "Soha 1", "desc": "Qo'llanilishi (8-12 so'z)"}},
    {{"icon": "🎓", "title": "Soha 2", "desc": "Qo'llanilishi (8-12 so'z)"}},
    {{"icon": "🚗", "title": "Soha 3", "desc": "Qo'llanilishi (8-12 so'z)"}},
    {{"icon": "💰", "title": "Soha 4", "desc": "Qo'llanilishi (8-12 so'z)"}}
  ],
  "conclusion_main": "Asosiy xulosa (2-3 so'z)",
  "conclusion_sub": "Davomi (2-3 so'z, masalan: bugun yaratiladi)",
  "conclusion_points": [
    "1-asosiy fikr (3-5 so'z)",
    "2-asosiy fikr (3-5 so'z)",
    "3-asosiy fikr (3-5 so'z)"
  ]
}}

MUHIM: Faqat JSON, boshqa matn yo'q. O'zbek tilida yozing."""

    response = requests.post(
        f"{GEMINI_URL}?key={GEMINI_API_KEY}",
        headers={"Content-Type": "application/json"},
        json={
            "contents": [{"parts": [{"text": prompt}]}],
            "generationConfig": {
                "temperature": 0.7,
                "maxOutputTokens": 2048,
                "responseMimeType": "application/json"
            }
        },
        timeout=60
    )

    if response.status_code != 200:
        raise Exception(f"Gemini API xato: {response.status_code} - {response.text[:200]}")

    data = response.json()
    text = data['candidates'][0]['content']['parts'][0]['text']
    return json.loads(text)


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Hex rangdan RGBColor yaratadi."""
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def add_rounded_rect(slide, x, y, w, h, fill_color, text=None, font_size=14,
                     font_color="FFFFFF", bold=False, italic=False, align=PP_ALIGN.LEFT):
    """Yumaloq burchakli to'rtburchak qo'shadi."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
    shape.line.fill.background()

    # Yumaloqlik darajasi
    shape.adjustments[0] = 0.1

    if text:
        tf = shape.text_frame
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = hex_to_rgb(font_color)
        run.font.name = "Calibri"
    return shape


def add_text_box(slide, x, y, w, h, text, font_size=14, color="FFFFFF",
                 bold=False, italic=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    """Matn qutisi qo'shadi."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = hex_to_rgb(color)
    run.font.name = font_name
    return tb


def create_bento_presentation(content: dict) -> bytes:
    """Bento Box uslubida prezentatsiya yaratadi."""
    # Ranglar (Apple uslubi)
    BG = "0A0A0A"
    CARD_DARK = "1A1A1C"
    CARD_GRAY = "232325"
    BLUE = "0A84FF"
    PURPLE = "BF5AF2"
    ORANGE = "FF9F0A"
    GREEN = "30D158"
    PINK = "FF375F"
    WHITE = "FFFFFF"
    GRAY = "8E8E93"
    LIGHT = "EBEBF5"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def set_bg(slide, color):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = hex_to_rgb(color)
        bg.line.fill.background()
        # Orqaga jo'natamiz
        spTree = bg._element.getparent()
        spTree.remove(bg._element)
        spTree.insert(2, bg._element)

    # ─────── SLAYD 1: TITUL ───────
    s1 = prs.slides.add_slide(blank)
    set_bg(s1, BG)

    # Asosiy karta
    add_rounded_rect(s1, 0.4, 0.4, 8.5, 4, CARD_DARK)
    add_text_box(s1, 0.8, 0.8, 6, 0.5, "AI", 14, BLUE, bold=True)
    add_text_box(s1, 0.8, 1.6, 8, 1.3, content['title_main'], 64, WHITE, bold=True)
    add_text_box(s1, 0.8, 2.7, 8, 1.3, content['title_sub'] + ".", 64, BLUE, bold=True)

    # Yon karta 1 (Kelajak)
    add_rounded_rect(s1, 9.1, 0.4, 3.8, 1.9, BLUE)
    add_text_box(s1, 9.4, 0.6, 1, 0.6, "✨", 32)
    add_text_box(s1, 9.4, 1.2, 3.4, 1, "Kelajak\nbugun", 22, WHITE, bold=True)

    # Yon karta 2 (Sana)
    add_rounded_rect(s1, 9.1, 2.5, 3.8, 1.9, CARD_GRAY)
    add_text_box(s1, 9.4, 2.7, 3.4, 1.2, "2026", 60, ORANGE, bold=True)
    add_text_box(s1, 9.4, 3.85, 3.4, 0.4, "Akademik yil", 12, GRAY)

    # Pastki katta karta
    add_rounded_rect(s1, 0.4, 4.5, 12.5, 2.6, CARD_DARK)
    add_text_box(s1, 0.8, 4.8, 12, 0.5,
                 "Texnologiya • Kelajak • Imkoniyat", 12, GRAY)
    add_text_box(s1, 0.8, 5.4, 8.5, 1.5,
                 f"Ushbu taqdimotda biz {content['title_main'].lower()} mavzusini batafsil o'rganamiz.",
                 18, LIGHT)
    add_text_box(s1, 11.8, 5.8, 0.8, 0.8, "→", 48, BLUE, bold=True, align=PP_ALIGN.CENTER)

    # ─────── SLAYD 2: KIRISH ───────
    s2 = prs.slides.add_slide(blank)
    set_bg(s2, BG)

    add_text_box(s2, 0.4, 0.3, 6, 0.4, "01 — Kirish", 12, GRAY)

    # Asosiy ta'rif karta
    add_rounded_rect(s2, 0.4, 0.9, 7, 6.2, CARD_DARK)
    add_text_box(s2, 0.8, 1.3, 6.3, 1, content['intro_question'], 48, WHITE, bold=True)

    # Kichik chiziq
    line = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0.8), Inches(2.5), Inches(6.3), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(BLUE)
    line.line.fill.background()

    add_text_box(s2, 0.8, 2.9, 6.3, 2, content['intro_definition'], 18, LIGHT)
    add_text_box(s2, 0.8, 5, 6.3, 1.5, content['intro_extra'], 14, GRAY)

    # Statistika 1
    add_rounded_rect(s2, 7.6, 0.9, 5.3, 3, BLUE)
    add_text_box(s2, 7.9, 1.1, 4.7, 1.6, content['stat1_value'], 88, WHITE, bold=True)
    add_text_box(s2, 7.9, 2.8, 4.7, 1, content['stat1_label'], 14, WHITE)

    # Statistika 2
    add_rounded_rect(s2, 7.6, 4.1, 5.3, 3, CARD_GRAY)
    add_text_box(s2, 7.9, 4.4, 4.7, 1.6, content['stat2_value'], 76, GREEN, bold=True)
    add_text_box(s2, 7.9, 5.95, 4.7, 1, content['stat2_label'], 14, LIGHT)

    # ─────── SLAYD 3: YO'NALISHLAR ───────
    s3 = prs.slides.add_slide(blank)
    set_bg(s3, BG)
    add_text_box(s3, 0.4, 0.3, 8, 0.4, "02 — Asosiy yo'nalishlar", 12, GRAY)

    dirs = content['directions']

    # 1-karta (katta, chap yuqori)
    add_rounded_rect(s3, 0.4, 0.9, 6.3, 3.2, BLUE)
    add_text_box(s3, 0.7, 1.1, 1, 0.4, "01", 12, WHITE)
    add_text_box(s3, 5.5, 1.1, 1, 1, dirs[0]['icon'], 36)
    add_text_box(s3, 0.7, 2.3, 5.5, 0.7, dirs[0]['title'], 28, WHITE, bold=True)
    add_text_box(s3, 0.7, 3.05, 5.5, 0.9, dirs[0]['desc'], 13, LIGHT)

    # 2-karta (o'ng yuqori)
    add_rounded_rect(s3, 6.9, 0.9, 6, 3.2, CARD_DARK)
    add_text_box(s3, 7.2, 1.1, 1, 0.4, "02", 12, GRAY)
    add_text_box(s3, 11.5, 1.1, 1, 1, dirs[1]['icon'], 36)
    add_text_box(s3, 7.2, 2.3, 5.5, 0.7, dirs[1]['title'], 28, PURPLE, bold=True)
    add_text_box(s3, 7.2, 3.05, 5.5, 0.9, dirs[1]['desc'], 13, LIGHT)

    # 3-karta (chap pastda)
    add_rounded_rect(s3, 0.4, 4.3, 6.3, 2.8, CARD_GRAY)
    add_text_box(s3, 0.7, 4.5, 1, 0.4, "03", 12, GRAY)
    add_text_box(s3, 5.5, 4.5, 1, 1, dirs[2]['icon'], 36)
    add_text_box(s3, 0.7, 5.5, 5.5, 0.7, dirs[2]['title'], 24, ORANGE, bold=True)
    add_text_box(s3, 0.7, 6.2, 5.5, 0.7, dirs[2]['desc'], 13, LIGHT)

    # 4-karta (o'ng pastda)
    add_rounded_rect(s3, 6.9, 4.3, 6, 2.8, PINK)
    add_text_box(s3, 7.2, 4.5, 1, 0.4, "04", 12, WHITE)
    add_text_box(s3, 11.5, 4.5, 1, 1, dirs[3]['icon'], 36)
    add_text_box(s3, 7.2, 5.5, 5.5, 0.7, dirs[3]['title'], 24, WHITE, bold=True)
    add_text_box(s3, 7.2, 6.2, 5.5, 0.7, dirs[3]['desc'], 13, LIGHT)

    # ─────── SLAYD 4: AMALIY QO'LLANISH ───────
    s4 = prs.slides.add_slide(blank)
    set_bg(s4, BG)
    add_text_box(s4, 0.4, 0.3, 8, 0.4, "03 — Hayotda qo'llanilishi", 12, GRAY)

    apps = content['applications']
    bgs = [GREEN, CARD_DARK, CARD_GRAY, PURPLE]
    title_colors = [WHITE, BLUE, ORANGE, WHITE]
    desc_colors = [LIGHT, GRAY, GRAY, LIGHT]

    for i, app in enumerate(apps):
        x = 0.4 + (i % 2) * 6.5
        y = 0.9 + (i // 2) * 3.15
        add_rounded_rect(s4, x, y, 6.4, 2.9, bgs[i])
        add_text_box(s4, x + 0.3, y + 0.3, 1, 1, app['icon'], 36)
        add_text_box(s4, x + 0.3, y + 1.5, 5.8, 0.7, app['title'], 26, title_colors[i], bold=True)
        add_text_box(s4, x + 0.3, y + 2.2, 5.8, 0.7, app['desc'], 13, desc_colors[i])

    # ─────── SLAYD 5: XULOSA ───────
    s5 = prs.slides.add_slide(blank)
    set_bg(s5, BG)
    add_text_box(s5, 0.4, 0.3, 6, 0.4, "04 — Xulosa", 12, GRAY)

    # Katta xulosa karta
    add_rounded_rect(s5, 0.4, 0.9, 12.5, 4, CARD_DARK)
    add_text_box(s5, 0.8, 1.4, 11.7, 1.5, content['conclusion_main'], 72, WHITE, bold=True)
    add_text_box(s5, 0.8, 2.8, 11.7, 1.5, content['conclusion_sub'] + ".", 72, BLUE, bold=True)

    # 3 ta kichik karta
    point_colors = [GREEN, ORANGE, PINK]
    points = content['conclusion_points']
    for i, point in enumerate(points[:3]):
        x = 0.4 + i * 4.2
        add_rounded_rect(s5, x, 5.1, 4, 1.4, CARD_GRAY)
        add_text_box(s5, x + 0.3, 5.25, 1, 0.3, f"0{i+1}", 11, point_colors[i])
        add_text_box(s5, x + 0.3, 5.6, 3.5, 0.8, point, 14, WHITE, bold=True)

    add_text_box(s5, 0.4, 6.7, 8, 0.4, "→  Rahmat e'tiboringiz uchun!", 14, BLUE, bold=True)
    add_text_box(s5, 8, 6.7, 4.9, 0.4, f"{content['title_main']} • 2026",
                 11, GRAY, align=PP_ALIGN.RIGHT)

    # Faylni byte'larga saqlash
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


# ─────────────────────────────────────────────────────────────────

# ─────────────────────────── MAJBURIY OBUNA ──────────────────────
# Kanal username — @ belgisi bilan yozing (masalan: "@edusave_kanal")
# Agar majburiy obuna kerak bo'lmasa, REQUIRED_CHANNEL = None qiling
REQUIRED_CHANNEL = "@paradox_sigma"   # ← Shu yerga o'z kanalingizni yozing
CHANNEL_LINK     = "https://t.me/paradox_sigma"  # ← Kanalga link

async def is_subscribed(bot, user_id: int) -> bool:
    """Foydalanuvchi kanalga obuna ekanligini tekshiradi."""
    if not REQUIRED_CHANNEL:
        return True
    try:
        member = await bot.get_chat_member(REQUIRED_CHANNEL, user_id)
        status = member.status
        logger.info(f"User {user_id} status in {REQUIRED_CHANNEL}: {status}")
        # 'left' va 'kicked' — obuna emas; qolgan barcha statuslar — obuna
        return status not in ("left", "kicked")
    except Exception as e:
        logger.error(f"❌ is_subscribed error for {user_id}: {e}")
        logger.error(f"   Bot {REQUIRED_CHANNEL} ga admin qilinganmi? Kanal username to'g'rimi?")
        # Xato chiqsa - obuna yo'q deb hisoblaymiz (xavfsizroq)
        return False

def subscribe_kb():
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("📢 Kanalga obuna bo'lish", url=CHANNEL_LINK)],
        [InlineKeyboardButton("✅ Tekshirish",            callback_data="check_sub")],
    ])

async def require_subscription(update: Update, ctx: ContextTypes.DEFAULT_TYPE) -> bool:
    """
    Foydalanuvchini tekshiradi. Obuna bo'lmasa - majburiy obuna xabarini chiqaradi.
    True qaytarsa - obuna bor, davom etish mumkin
    False qaytarsa - obuna yo'q, to'xtatish kerak
    """
    uid = update.effective_user.id
    if await is_subscribed(ctx.bot, uid):
        return True

    text = (
        "🔒 <b>Botdan foydalanish uchun kanalimizga obuna bo'ling!</b>\n\n"
        "1. Quyidagi tugma orqali kanalga o'ting\n"
        "2. Obuna bo'ling\n"
        "3. <b>✅ Tekshirish</b> tugmasini bosing"
    )
    if update.callback_query:
        try:
            await update.callback_query.answer("🔒 Avval kanalga obuna bo'ling!", show_alert=True)
            await update.callback_query.edit_message_text(
                text, reply_markup=subscribe_kb(), parse_mode="HTML"
            )
        except Exception:
            await update.callback_query.message.reply_text(
                text, reply_markup=subscribe_kb(), parse_mode="HTML"
            )
    elif update.message:
        await update.message.reply_text(text, reply_markup=subscribe_kb(), parse_mode="HTML")
    return False

# ─────────────────────────── DATABASE ────────────────────────────

def get_conn():
    conn = sqlite3.connect(DB_PATH)
    conn.execute("PRAGMA foreign_keys = ON")
    return conn

def init_db():
    with get_conn() as conn:
        conn.executescript('''
            CREATE TABLE IF NOT EXISTS users (
                user_id      INTEGER PRIMARY KEY,
                username     TEXT,
                first_name   TEXT,
                first_seen   TEXT DEFAULT (datetime('now')),
                last_seen    TEXT DEFAULT (datetime('now'))
            );
            CREATE TABLE IF NOT EXISTS categories (
                id         INTEGER PRIMARY KEY AUTOINCREMENT,
                user_id    INTEGER NOT NULL,
                name       TEXT    NOT NULL,
                emoji      TEXT    DEFAULT '📁',
                created_at TEXT    DEFAULT (datetime('now')),
                UNIQUE(user_id, name)
            );
            CREATE TABLE IF NOT EXISTS items (
                id           INTEGER PRIMARY KEY AUTOINCREMENT,
                user_id      INTEGER NOT NULL,
                category_id  INTEGER NOT NULL,
                msg_type     TEXT    NOT NULL,
                text_content TEXT,
                file_id      TEXT,
                file_name    TEXT,
                caption      TEXT,
                tags         TEXT    DEFAULT '',
                is_favorite  INTEGER DEFAULT 0,
                created_at   TEXT    DEFAULT (datetime('now')),
                FOREIGN KEY (category_id) REFERENCES categories(id) ON DELETE CASCADE
            );
        ''')

# ─────────────────────────── USER TRACKING ───────────────────────

def track_user(user_id: int, username: str = None, first_name: str = None):
    """Foydalanuvchini bazaga qo'shadi yoki yangilaydi."""
    with get_conn() as conn:
        conn.execute(
            "INSERT INTO users (user_id, username, first_name) VALUES (?,?,?) "
            "ON CONFLICT(user_id) DO UPDATE SET "
            "  username=excluded.username, "
            "  first_name=excluded.first_name, "
            "  last_seen=datetime('now')",
            (user_id, username, first_name)
        )

def get_admin_stats():
    """Admin uchun statistika."""
    with get_conn() as conn:
        total_users    = conn.execute("SELECT COUNT(*) FROM users").fetchone()[0]
        today_users    = conn.execute("SELECT COUNT(*) FROM users WHERE date(first_seen)=date('now')").fetchone()[0]
        week_users     = conn.execute("SELECT COUNT(*) FROM users WHERE first_seen >= datetime('now','-7 days')").fetchone()[0]
        active_today   = conn.execute("SELECT COUNT(*) FROM users WHERE date(last_seen)=date('now')").fetchone()[0]
        total_cats     = conn.execute("SELECT COUNT(*) FROM categories").fetchone()[0]
        total_items    = conn.execute("SELECT COUNT(*) FROM items").fetchone()[0]
        top_users = conn.execute(
            "SELECT u.user_id, u.first_name, u.username, "
            "       (SELECT COUNT(*) FROM items WHERE user_id=u.user_id) as item_cnt "
            "FROM users u "
            "ORDER BY item_cnt DESC LIMIT 5"
        ).fetchall()
    return {
        'total_users':  total_users,
        'today_users':  today_users,
        'week_users':   week_users,
        'active_today': active_today,
        'total_cats':   total_cats,
        'total_items':  total_items,
        'top_users':    top_users,
    }

def get_recent_users(limit=10):
    """So'nggi foydalanuvchilar."""
    with get_conn() as conn:
        return conn.execute(
            "SELECT user_id, first_name, username, first_seen FROM users "
            "ORDER BY first_seen DESC LIMIT ?",
            (limit,)
        ).fetchall()

# ─────────────────────────── DB HELPERS ──────────────────────────

def get_categories(user_id):
    with get_conn() as conn:
        return conn.execute(
            "SELECT id, name, emoji FROM categories WHERE user_id=? ORDER BY name",
            (user_id,)
        ).fetchall()

def add_category(user_id, name, emoji='📁'):
    try:
        with get_conn() as conn:
            conn.execute(
                "INSERT INTO categories (user_id, name, emoji) VALUES (?,?,?)",
                (user_id, name, emoji)
            )
        return True
    except sqlite3.IntegrityError:
        return False

def delete_category(user_id, cat_id):
    with get_conn() as conn:
        conn.execute("DELETE FROM categories WHERE id=? AND user_id=?", (cat_id, user_id))

def get_category(user_id, cat_id):
    with get_conn() as conn:
        return conn.execute(
            "SELECT id, name, emoji FROM categories WHERE id=? AND user_id=?",
            (cat_id, user_id)
        ).fetchone()

def save_item(user_id, cat_id, msg_type, text=None, file_id=None,
              file_name=None, caption=None, tags=''):
    with get_conn() as conn:
        conn.execute(
            "INSERT INTO items "
            "(user_id,category_id,msg_type,text_content,file_id,file_name,caption,tags)"
            " VALUES (?,?,?,?,?,?,?,?)",
            (user_id, cat_id, msg_type, text, file_id, file_name, caption, tags)
        )

def get_items(user_id, cat_id, page=0, search=None):
    offset = page * ITEMS_PER_PAGE
    with get_conn() as conn:
        if search:
            q = f'%{search}%'
            rows = conn.execute(
                "SELECT id,msg_type,text_content,file_name,caption,tags,is_favorite,created_at"
                " FROM items WHERE user_id=? AND category_id=?"
                " AND (text_content LIKE ? OR caption LIKE ? OR tags LIKE ? OR file_name LIKE ?)"
                " ORDER BY is_favorite DESC, created_at DESC LIMIT ? OFFSET ?",
                (user_id, cat_id, q, q, q, q, ITEMS_PER_PAGE, offset)
            ).fetchall()
            total = conn.execute(
                "SELECT COUNT(*) FROM items WHERE user_id=? AND category_id=?"
                " AND (text_content LIKE ? OR caption LIKE ? OR tags LIKE ? OR file_name LIKE ?)",
                (user_id, cat_id, q, q, q, q)
            ).fetchone()[0]
        else:
            rows = conn.execute(
                "SELECT id,msg_type,text_content,file_name,caption,tags,is_favorite,created_at"
                " FROM items WHERE user_id=? AND category_id=?"
                " ORDER BY is_favorite DESC, created_at DESC LIMIT ? OFFSET ?",
                (user_id, cat_id, ITEMS_PER_PAGE, offset)
            ).fetchall()
            total = conn.execute(
                "SELECT COUNT(*) FROM items WHERE user_id=? AND category_id=?",
                (user_id, cat_id)
            ).fetchone()[0]
    return rows, total

def get_item(user_id, item_id):
    with get_conn() as conn:
        return conn.execute(
            "SELECT id,msg_type,text_content,file_id,file_name,caption,tags,is_favorite"
            " FROM items WHERE id=? AND user_id=?",
            (item_id, user_id)
        ).fetchone()

def get_all_file_items(user_id, cat_id):
    with get_conn() as conn:
        return conn.execute(
            "SELECT id,msg_type,file_id,file_name,caption"
            " FROM items WHERE user_id=? AND category_id=? AND file_id IS NOT NULL"
            " ORDER BY created_at",
            (user_id, cat_id)
        ).fetchall()

def toggle_favorite(user_id, item_id):
    with get_conn() as conn:
        conn.execute(
            "UPDATE items SET is_favorite = 1 - is_favorite WHERE id=? AND user_id=?",
            (item_id, user_id)
        )

def delete_item(user_id, item_id):
    with get_conn() as conn:
        conn.execute("DELETE FROM items WHERE id=? AND user_id=?", (item_id, user_id))

def search_all_items(user_id, query):
    q = f'%{query}%'
    with get_conn() as conn:
        return conn.execute(
            "SELECT i.id,i.msg_type,i.text_content,i.file_name,i.caption,"
            "       i.is_favorite,i.created_at,c.name,c.emoji,i.category_id"
            " FROM items i JOIN categories c ON i.category_id=c.id"
            " WHERE i.user_id=?"
            " AND (i.text_content LIKE ? OR i.caption LIKE ? OR i.tags LIKE ? OR i.file_name LIKE ?)"
            " ORDER BY i.is_favorite DESC, i.created_at DESC LIMIT 15",
            (user_id, q, q, q, q)
        ).fetchall()

def get_favorites(user_id):
    with get_conn() as conn:
        return conn.execute(
            "SELECT i.id,i.msg_type,i.text_content,i.file_name,i.caption,"
            "       i.created_at,c.name,c.emoji,i.category_id"
            " FROM items i JOIN categories c ON i.category_id=c.id"
            " WHERE i.user_id=? AND i.is_favorite=1"
            " ORDER BY i.created_at DESC LIMIT 20",
            (user_id,)
        ).fetchall()

def get_stats(user_id):
    with get_conn() as conn:
        total_cats  = conn.execute("SELECT COUNT(*) FROM categories WHERE user_id=?", (user_id,)).fetchone()[0]
        total_items = conn.execute("SELECT COUNT(*) FROM items WHERE user_id=?", (user_id,)).fetchone()[0]
        favorites   = conn.execute("SELECT COUNT(*) FROM items WHERE user_id=? AND is_favorite=1", (user_id,)).fetchone()[0]
        top_cats    = conn.execute(
            "SELECT c.name,c.emoji,COUNT(i.id) cnt FROM categories c"
            " LEFT JOIN items i ON c.id=i.category_id WHERE c.user_id=?"
            " GROUP BY c.id ORDER BY cnt DESC LIMIT 5",
            (user_id,)
        ).fetchall()
    return total_cats, total_items, favorites, top_cats

def cat_item_count(user_id, cat_id):
    with get_conn() as conn:
        return conn.execute(
            "SELECT COUNT(*) FROM items WHERE user_id=? AND category_id=?",
            (user_id, cat_id)
        ).fetchone()[0]

# ─────────────────────────── ICONS ───────────────────────────────

TYPE_ICON = {
    "text": "📝", "photo": "🖼", "document": "📄",
    "video": "🎬", "audio": "🎵", "voice": "🎤", "sticker": "😊"
}
TYPE_NAME = {
    "text": "Matn", "photo": "Rasm", "document": "Hujjat",
    "video": "Video", "audio": "Audio", "voice": "Ovozli xabar"
}
COMPRESSIBLE = {"photo", "document", "audio", "video", "voice"}

def item_preview(row):
    _, msg_type, text, fname, cap, *_ = row
    raw = text or cap or fname or "📎 Fayl"
    return TYPE_ICON.get(msg_type, "📎") + " " + (raw[:35] + "…" if len(raw) > 35 else raw)

def human_size(nbytes: int) -> str:
    for unit in ("B", "KB", "MB", "GB"):
        if nbytes < 1024:
            return f"{nbytes:.1f} {unit}"
        nbytes /= 1024
    return f"{nbytes:.1f} TB"

# ─────────────────────────── COMPRESSION ─────────────────────────

QUALITY_LEVELS = {
    "yuqori":  85,   # engil siqish  ~30% kichrayadi
    "urtacha": 60,   # o'rtacha      ~55% kichrayadi
    "kuchli":  35,   # kuchli siqish ~75% kichrayadi
}

async def compress_image_bytes(data: bytes, quality: int) -> bytes:
    img = Image.open(io.BytesIO(data))
    if img.mode in ("RGBA", "P"):
        img = img.convert("RGB")
    out = io.BytesIO()
    img.save(out, format="JPEG", quality=quality, optimize=True)
    return out.getvalue()

async def compress_generic_bytes(data: bytes, filename: str) -> tuple:
    out = io.BytesIO()
    zname = (filename or "fayl") + ".zip"
    with zipfile.ZipFile(out, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=9) as zf:
        zf.writestr(filename or "fayl", data)
    return out.getvalue(), zname

async def download_file(bot, file_id: str) -> bytes:
    tg_file = await bot.get_file(file_id)
    buf = io.BytesIO()
    await tg_file.download_to_memory(buf)
    return buf.getvalue()

async def build_zip_for_category(bot, user_id: int, cat_id: int) -> io.BytesIO | None:
    file_items = get_all_file_items(user_id, cat_id)
    if not file_items:
        return None
    zip_buf = io.BytesIO()
    counters = {}
    with zipfile.ZipFile(zip_buf, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=6) as zf:
        for item_id, msg_type, file_id, file_name, caption in file_items:
            try:
                data = await download_file(bot, file_id)
            except Exception as e:
                logger.warning(f"ZIP skip {file_id}: {e}")
                continue
            if file_name:
                arc_name = file_name
            else:
                ext = {"photo": "jpg", "video": "mp4", "audio": "mp3", "voice": "ogg"}.get(msg_type, "bin")
                arc_name = f"{msg_type}_{item_id}.{ext}"
            if arc_name in counters:
                counters[arc_name] += 1
                base, _, ext2 = arc_name.rpartition(".")
                arc_name = f"{base}_{counters[arc_name]}.{ext2}" if ext2 else f"{arc_name}_{counters[arc_name]}"
            else:
                counters[arc_name] = 0
            zf.writestr(arc_name, data)
    zip_buf.seek(0)
    return zip_buf

# ─────────────────────────── KEYBOARDS ───────────────────────────

def main_menu_kb():
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("📚 Kategoriyalar", callback_data="cat_list"),
         InlineKeyboardButton("🔍 Qidirish",      callback_data="search_start")],
        [InlineKeyboardButton("➕ Kategoriya",    callback_data="cat_add"),
         InlineKeyboardButton("⭐ Sevimlilar",    callback_data="favorites")],
        [InlineKeyboardButton("📊 Statistika",    callback_data="stats"),
         InlineKeyboardButton("📦 ZIP qilish",    callback_data="zip_mode")],
        [InlineKeyboardButton("🎨 AI Prezentatsiya",  callback_data="ai_preso")],
    ])

def back_btn(target="main_menu"):
    return InlineKeyboardButton("🔙 Orqaga", callback_data=target)

def categories_kb(user_id):
    cats = get_categories(user_id)
    rows = []
    for cat_id, name, emoji in cats:
        cnt = cat_item_count(user_id, cat_id)
        rows.append([InlineKeyboardButton(
            f"{emoji} {name}  ({cnt})",
            callback_data=f"cat_open:{cat_id}:0"
        )])
    rows.append([back_btn("main_menu")])
    return InlineKeyboardMarkup(rows)

def items_list_kb(user_id, cat_id, page, total, items):
    buttons = []
    for row in items:
        item_id = row[0]
        star = "⭐ " if row[6] else ""
        buttons.append([InlineKeyboardButton(
            star + item_preview(row),
            callback_data=f"item_view:{item_id}:{cat_id}:{page}"
        )])
    nav = []
    total_pages = max(1, (total + ITEMS_PER_PAGE - 1) // ITEMS_PER_PAGE)
    if page > 0:
        nav.append(InlineKeyboardButton("◀️", callback_data=f"cat_open:{cat_id}:{page-1}"))
    nav.append(InlineKeyboardButton(f"{page+1}/{total_pages}", callback_data="noop"))
    if (page + 1) * ITEMS_PER_PAGE < total:
        nav.append(InlineKeyboardButton("▶️", callback_data=f"cat_open:{cat_id}:{page+1}"))
    if nav:
        buttons.append(nav)
    buttons.append([
        InlineKeyboardButton("🔍 Izlash",     callback_data=f"cat_search:{cat_id}"),
        InlineKeyboardButton("📦 ZIP arxiv",  callback_data=f"zip_cat:{cat_id}"),
    ])
    buttons.append([
        InlineKeyboardButton("🗑 Kategoriyani o'chir", callback_data=f"cat_del_confirm:{cat_id}"),
        back_btn("cat_list"),
    ])
    return InlineKeyboardMarkup(buttons)

def item_actions_kb(item_id, cat_id, page, is_fav, msg_type):
    star_label = "💛 Sevimlilardan chiqar" if is_fav else "⭐ Sevimliga qo'sh"
    rows = [
        [InlineKeyboardButton(star_label,  callback_data=f"item_fav:{item_id}:{cat_id}:{page}"),
         InlineKeyboardButton("🗑 O'chir", callback_data=f"item_del:{item_id}:{cat_id}:{page}")],
    ]
    if msg_type in COMPRESSIBLE:
        rows.append([
            InlineKeyboardButton("🗜 Engil",   callback_data=f"compress:{item_id}:{cat_id}:{page}:yuqori"),
            InlineKeyboardButton("🗜 O'rtacha", callback_data=f"compress:{item_id}:{cat_id}:{page}:urtacha"),
            InlineKeyboardButton("🗜 Kuchli",   callback_data=f"compress:{item_id}:{cat_id}:{page}:kuchli"),
        ])
    rows.append([back_btn(f"cat_open:{cat_id}:{page}")])
    return InlineKeyboardMarkup(rows)

# ──────────────────────────── HANDLERS ───────────────────────────

async def cmd_start(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    logger.info(f"🟢 /start dan: user_id={update.effective_user.id}, name={update.effective_user.first_name}")
    # Foydalanuvchini ro'yxatga olamiz
    u = update.effective_user
    track_user(u.id, u.username, u.first_name)
    if not await require_subscription(update, ctx):
        logger.info(f"🔴 Obuna emas, to'xtatildi: {update.effective_user.id}")
        return
    logger.info(f"🟢 Obuna tasdiqlandi: {update.effective_user.id}")
    name = update.effective_user.first_name
    text = (
        f"👋 Salom, <b>{name}</b>!\n\n"
        "📚 <b>EduSave Bot</b> — o'quv materiallaringizni tartibli saqlang!\n\n"
        "<b>Imkoniyatlar:</b>\n"
        "• Cheksiz kategoriya (📐 Matematika, 📖 Ona tili…)\n"
        "• Matn, rasm, hujjat, video, audio saqlash\n"
        "• 🗜 Faylni siqish — 3 daraja: engil / o'rtacha / kuchli\n"
        "• 📦 Kategoriyani ZIP arxivga yuklash\n"
        "• 🔍 Global va kategoriyada qidirish\n"
        "• ⭐ Sevimlilar va 📊 Statistika\n\n"
        "<b>Qanday saqlash:</b>\n"
        "Istalgan xabar yuboring → kategoriyani tanlang → tayyor! ✅"
    )
    await update.message.reply_text(text, reply_markup=main_menu_kb(), parse_mode="HTML")

# ─────────────────────────── ADMIN PANEL ─────────────────────────

def admin_menu_kb():
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("👥 Foydalanuvchilar", callback_data="admin_users"),
         InlineKeyboardButton("📊 Statistika",       callback_data="admin_stats")],
        [InlineKeyboardButton("🆕 So'nggi qo'shilganlar", callback_data="admin_recent")],
        [InlineKeyboardButton("📢 Barchaga xabar yuborish", callback_data="admin_broadcast")],
        [InlineKeyboardButton("🔙 Asosiy menyu", callback_data="main_menu")],
    ])

async def cmd_admin(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    """Admin paneli — faqat ADMIN_ID ko'rishi mumkin."""
    uid = update.effective_user.id
    if uid != ADMIN_ID:
        await update.message.reply_text("❌ Sizda admin huquqlari yo'q")
        return
    stats = get_admin_stats()
    text = (
        "👨‍💼 <b>Admin Paneli</b>\n\n"
        f"👥 Jami foydalanuvchilar: <b>{stats['total_users']}</b>\n"
        f"🆕 Bugun qo'shilganlar: <b>{stats['today_users']}</b>\n"
        f"📅 Bu hafta qo'shilganlar: <b>{stats['week_users']}</b>\n"
        f"🟢 Bugun faol: <b>{stats['active_today']}</b>\n\n"
        f"📁 Jami kategoriyalar: <b>{stats['total_cats']}</b>\n"
        f"📦 Jami materiallar: <b>{stats['total_items']}</b>"
    )
    await update.message.reply_text(text, reply_markup=admin_menu_kb(), parse_mode="HTML")

async def cb_admin_stats(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    await q.answer()
    if q.from_user.id != ADMIN_ID:
        return
    stats = get_admin_stats()
    lines = [
        "📊 <b>To'liq statistika</b>\n",
        f"👥 Jami foydalanuvchilar: <b>{stats['total_users']}</b>",
        f"🆕 Bugun qo'shilganlar: <b>{stats['today_users']}</b>",
        f"📅 Bu hafta qo'shilganlar: <b>{stats['week_users']}</b>",
        f"🟢 Bugun faol: <b>{stats['active_today']}</b>",
        f"📁 Jami kategoriyalar: <b>{stats['total_cats']}</b>",
        f"📦 Jami materiallar: <b>{stats['total_items']}</b>",
    ]
    if stats['top_users']:
        lines.append("\n🏆 <b>Top 5 foydalanuvchi (materiallar bo'yicha):</b>")
        for i, (uid, fname, uname, cnt) in enumerate(stats['top_users'], 1):
            display = fname or uname or f"User {uid}"
            lines.append(f"  {i}. {display} — {cnt} ta")
    kb = InlineKeyboardMarkup([[InlineKeyboardButton("🔙 Orqaga", callback_data="admin_back")]])
    await q.edit_message_text("\n".join(lines), reply_markup=kb, parse_mode="HTML")

async def cb_admin_users(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    await q.answer()
    if q.from_user.id != ADMIN_ID:
        return
    with get_conn() as conn:
        users = conn.execute(
            "SELECT user_id, first_name, username, first_seen, last_seen FROM users "
            "ORDER BY last_seen DESC LIMIT 20"
        ).fetchall()
    if not users:
        kb = InlineKeyboardMarkup([[InlineKeyboardButton("🔙 Orqaga", callback_data="admin_back")]])
        await q.edit_message_text("👥 Foydalanuvchilar yo'q", reply_markup=kb)
        return
    lines = [f"👥 <b>So'nggi 20 foydalanuvchi (faollik bo'yicha):</b>\n"]
    for uid, fname, uname, first_seen, last_seen in users:
        display = fname or "Noma'lum"
        username_str = f" (@{uname})" if uname else ""
        lines.append(f"• <b>{display}</b>{username_str}\n  └ ID: <code>{uid}</code> | Oxirgi: {last_seen[:16]}")
    kb = InlineKeyboardMarkup([[InlineKeyboardButton("🔙 Orqaga", callback_data="admin_back")]])
    await q.edit_message_text("\n".join(lines), reply_markup=kb, parse_mode="HTML")

async def cb_admin_recent(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    await q.answer()
    if q.from_user.id != ADMIN_ID:
        return
    users = get_recent_users(15)
    if not users:
        kb = InlineKeyboardMarkup([[InlineKeyboardButton("🔙 Orqaga", callback_data="admin_back")]])
        await q.edit_message_text("🆕 Yangi foydalanuvchilar yo'q", reply_markup=kb)
        return
    lines = [f"🆕 <b>So'nggi 15 yangi foydalanuvchi:</b>\n"]
    for uid, fname, uname, first_seen in users:
        display = fname or "Noma'lum"
        username_str = f" (@{uname})" if uname else ""
        lines.append(f"• <b>{display}</b>{username_str}\n  └ Qo'shildi: {first_seen[:16]}")
    kb = InlineKeyboardMarkup([[InlineKeyboardButton("🔙 Orqaga", callback_data="admin_back")]])
    await q.edit_message_text("\n".join(lines), reply_markup=kb, parse_mode="HTML")

async def cb_admin_back(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    """Admin paneliga qaytish."""
    q = update.callback_query
    await q.answer()
    if q.from_user.id != ADMIN_ID:
        return
    stats = get_admin_stats()
    text = (
        "👨‍💼 <b>Admin Paneli</b>\n\n"
        f"👥 Jami foydalanuvchilar: <b>{stats['total_users']}</b>\n"
        f"🆕 Bugun qo'shilganlar: <b>{stats['today_users']}</b>\n"
        f"📅 Bu hafta qo'shilganlar: <b>{stats['week_users']}</b>\n"
        f"🟢 Bugun faol: <b>{stats['active_today']}</b>\n\n"
        f"📁 Jami kategoriyalar: <b>{stats['total_cats']}</b>\n"
        f"📦 Jami materiallar: <b>{stats['total_items']}</b>"
    )
    await q.edit_message_text(text, reply_markup=admin_menu_kb(), parse_mode="HTML")

async def cb_admin_broadcast(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    """Hamma foydalanuvchilarga xabar yuborish."""
    q = update.callback_query
    await q.answer()
    if q.from_user.id != ADMIN_ID:
        return
    ctx.user_data['state'] = 'broadcast'
    kb = InlineKeyboardMarkup([[InlineKeyboardButton("❌ Bekor qilish", callback_data="admin_back")]])
    await q.edit_message_text(
        "📢 <b>Barchaga xabar yuborish</b>\n\n"
        "Yubormoqchi bo'lgan xabaringizni yozing.\n"
        "Format saqlanadi (bold, italic, link va h.k.).\n\n"
        "⚠️ Diqqat: barcha foydalanuvchilarga yuboriladi!",
        reply_markup=kb, parse_mode="HTML"
    )

# ─────────────────────────────────────────────────────────────────
# ─────────────────────── 🎨 AI PREZENTATSIYA ──────────────────────

async def cb_ai_preso(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    """AI Prezentatsiya yaratish boshlash."""
    q = update.callback_query
    await q.answer()
    uid = q.from_user.id

    # Hozircha faqat admin uchun
    if uid != ADMIN_ID:
        kb = InlineKeyboardMarkup([[InlineKeyboardButton("🔙 Asosiy menyu", callback_data="main_menu")]])
        await q.edit_message_text(
            "🎨 <b>AI Prezentatsiya yaratish</b>\n\n"
            "🚧 Bu funksiya hozircha <b>test rejimida</b>.\n"
            "Tez orada hammaga ochiladi!\n\n"
            "📌 Funksiya tayyor bo'lgach, bot orqali bildirishnoma olasiz.",
            reply_markup=kb, parse_mode="HTML"
        )
        return

    # Admin uchun
    ctx.user_data['state'] = 'ai_preso_topic'
    kb = InlineKeyboardMarkup([[InlineKeyboardButton("❌ Bekor", callback_data="main_menu")]])
    await q.edit_message_text(
        "🎨 <b>AI Prezentatsiya yaratish</b> [TEST]\n\n"
        "Mavzuni yozing, men 30-60 soniyada zamonaviy "
        "Bento Box uslubidagi prezentatsiya yarataman.\n\n"
        "<b>Misollar:</b>\n"
        "• Sun'iy intellekt\n"
        "• Ekologiya muammolari\n"
        "• Sog'lom turmush tarzi\n"
        "• Kvant kompyuterlar\n\n"
        "✍️ Mavzu yozing:",
        reply_markup=kb, parse_mode="HTML"
    )


async def handle_ai_preso(update: Update, ctx: ContextTypes.DEFAULT_TYPE, topic: str):
    """Prezentatsiya yaratadi va yuboradi."""
    msg = update.message
    uid = update.effective_user.id

    # Status xabari
    status = await msg.reply_text(
        "🎨 <b>Prezentatsiya yaratilmoqda...</b>\n\n"
        "⏳ Gemini AI matn yaratmoqda (15-30 soniya)",
        parse_mode="HTML"
    )

    try:
        # 1. Gemini'dan matn olish
        logger.info(f"🎨 AI Preso boshlandi: user={uid}, topic={topic}")
        content = generate_presentation_content(topic)
        logger.info(f"🎨 Gemini javob keldi: title={content.get('title_main')}")

        await status.edit_text(
            "🎨 <b>Prezentatsiya yaratilmoqda...</b>\n\n"
            "✅ Matn tayyor\n"
            "⏳ Slaydlar dizaynda joylanmoqda...",
            parse_mode="HTML"
        )

        # 2. Bento prezentatsiya yaratish
        pptx_data = create_bento_presentation(content)
        logger.info(f"🎨 PPTX yaratildi: {len(pptx_data)} bayt")

        # 3. Foydalanuvchiga yuborish
        await status.edit_text(
            "🎨 <b>Prezentatsiya yaratilmoqda...</b>\n\n"
            "✅ Matn tayyor\n"
            "✅ Dizayn tayyor\n"
            "📤 Yuborilmoqda...",
            parse_mode="HTML"
        )

        filename = f"{topic[:30].replace('/', '_')}.pptx"
        await msg.reply_document(
            io.BytesIO(pptx_data),
            filename=filename,
            caption=(
                f"🎨 <b>Prezentatsiya tayyor!</b>\n\n"
                f"📌 Mavzu: <b>{topic}</b>\n"
                f"📊 Slaydlar: <b>5 ta</b>\n"
                f"🎭 Uslub: <b>Bento Box (zamonaviy)</b>\n\n"
                f"💡 PowerPoint yoki Google Slides'da oching!"
            ),
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup([[
                InlineKeyboardButton("🏠 Asosiy menyu", callback_data="main_menu")
            ]])
        )

        await status.delete()
        ctx.user_data.clear()

    except Exception as e:
        logger.error(f"AI Preso xato: {e}", exc_info=True)
        await status.edit_text(
            f"❌ <b>Xato yuz berdi</b>\n\n"
            f"Sabab: <code>{str(e)[:200]}</code>\n\n"
            f"Iltimos, qayta urinib ko'ring yoki boshqa mavzu yozing.",
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup([[
                InlineKeyboardButton("🔙 Asosiy menyu", callback_data="main_menu")
            ]])
        )
        ctx.user_data.clear()

# ─────────────────────────────────────────────────────────────────

async def cb_main_menu(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    await q.answer()
    ctx.user_data.clear()
    await q.edit_message_text("🏠 <b>Asosiy menyu</b>", reply_markup=main_menu_kb(), parse_mode="HTML")

async def cb_cat_list(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    await q.answer()
    uid = q.from_user.id
    cats = get_categories(uid)
    text = (
        f"📚 <b>Kategoriyalar</b> ({len(cats)} ta)"
        if cats else
        "📂 Hali kategoriya yo'q.\n\n➕ Kategoriya qo'shing!"
    )
    await q.edit_message_text(text, reply_markup=categories_kb(uid), parse_mode="HTML")

async def cb_cat_add(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    await q.answer()
    ctx.user_data['state'] = 'adding_category'
    kb = InlineKeyboardMarkup([[InlineKeyboardButton("❌ Bekor", callback_data="cat_list")]])
    await q.edit_message_text(
        "➕ <b>Yangi kategoriya</b>\n\nKategoriya nomini yuboring. Emoji ham qo'shsangiz bo'ladi:\n\n"
        "<code>📐 Matematika</code>\n<code>📖 Ona tili</code>\n"
        "<code>🧪 Kimyo</code>\n<code>🌍 Geografiya</code>\n<code>💻 Informatika</code>",
        reply_markup=kb, parse_mode="HTML"
    )

async def cb_cat_open(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    await q.answer()
    uid = q.from_user.id
    _, cat_id, page = q.data.split(":")
    cat_id, page = int(cat_id), int(page)
    cat = get_category(uid, cat_id)
    if not cat:
        await q.edit_message_text("❌ Kategoriya topilmadi", reply_markup=main_menu_kb())
        return
    _, name, emoji = cat
    items, total = get_items(uid, cat_id, page=page)
    text = (
        f"{emoji} <b>{name}</b>\n📦 Jami: {total} ta material"
        if total else
        f"{emoji} <b>{name}</b>\n\n📭 Hali material yo'q.\nXabar yuboring va shu kategoriyani tanlang!"
    )
    await q.edit_message_text(text, reply_markup=items_list_kb(uid, cat_id, page, total, items), parse_mode="HTML")

async def cb_cat_del_confirm(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    await q.answer()
    uid = q.from_user.id
    cat_id = int(q.data.split(":")[1])
    cat = get_category(uid, cat_id)
    if not cat:
        return
    _, name, emoji = cat
    cnt = cat_item_count(uid, cat_id)
    kb = InlineKeyboardMarkup([[
        InlineKeyboardButton("✅ Ha, o'chir", callback_data=f"cat_del:{cat_id}"),
        InlineKeyboardButton("❌ Bekor",      callback_data=f"cat_open:{cat_id}:0"),
    ]])
    await q.edit_message_text(
        f"⚠️ <b>{emoji} {name}</b> kategoriyasini o'chirmoqchimisiz?\n"
        f"Ichidagi <b>{cnt}</b> ta material ham o'chib ketadi!",
        reply_markup=kb, parse_mode="HTML"
    )

async def cb_cat_del(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    uid = q.from_user.id
    cat_id = int(q.data.split(":")[1])
    delete_category(uid, cat_id)
    await q.answer("🗑 Kategoriya o'chirildi", show_alert=True)
    cats = get_categories(uid)
    await q.edit_message_text(
        f"📚 <b>Kategoriyalar</b> ({len(cats)} ta)",
        reply_markup=categories_kb(uid), parse_mode="HTML"
    )

async def cb_cat_search(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    await q.answer()
    cat_id = int(q.data.split(":")[1])
    ctx.user_data['state']      = 'searching_cat'
    ctx.user_data['search_cat'] = cat_id
    kb = InlineKeyboardMarkup([[InlineKeyboardButton("❌ Bekor", callback_data=f"cat_open:{cat_id}:0")]])
    await q.edit_message_text("🔍 <b>Kategoriyada qidirish</b>\n\nSo'zni yuboring:", reply_markup=kb, parse_mode="HTML")

async def cb_search_start(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    await q.answer()
    ctx.user_data['state'] = 'searching_all'
    kb = InlineKeyboardMarkup([[InlineKeyboardButton("❌ Bekor", callback_data="main_menu")]])
    await q.edit_message_text(
        "🔍 <b>Global qidirish</b>\n\nBarcha kategoriyalarda qidiradi.\nSo'zni yuboring:",
        reply_markup=kb, parse_mode="HTML"
    )

async def cb_item_view(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    await q.answer()
    uid  = q.from_user.id
    _, item_id, cat_id, page = q.data.split(":")
    item_id, cat_id, page = int(item_id), int(cat_id), int(page)
    row = get_item(uid, item_id)
    if not row:
        await q.answer("❌ Topilmadi", show_alert=True)
        return
    _, msg_type, text, file_id, file_name, caption, tags, is_fav = row
    kb = item_actions_kb(item_id, cat_id, page, is_fav, msg_type)
    tag_line = f"\n\n🏷 <i>{tags}</i>" if tags else ""
    try:
        if msg_type == "text":
            await q.edit_message_text(f"📝 <b>Matn</b>{tag_line}\n\n{text}", reply_markup=kb, parse_mode="HTML")
        elif msg_type == "photo":
            await q.message.reply_photo(file_id, caption=f"{caption or ''}{tag_line}", reply_markup=kb, parse_mode="HTML")
        elif msg_type == "document":
            await q.message.reply_document(file_id, caption=f"📄 {file_name or ''}\n{caption or ''}{tag_line}", reply_markup=kb, parse_mode="HTML")
        elif msg_type == "video":
            await q.message.reply_video(file_id, caption=f"{caption or ''}{tag_line}", reply_markup=kb, parse_mode="HTML")
        elif msg_type == "audio":
            await q.message.reply_audio(file_id, caption=f"{caption or ''}{tag_line}", reply_markup=kb, parse_mode="HTML")
        elif msg_type == "voice":
            await q.message.reply_voice(file_id, caption=f"{caption or ''}{tag_line}", reply_markup=kb, parse_mode="HTML")
    except Exception as e:
        logger.warning(f"item_view error: {e}")
        await q.answer("❌ Faylni ko'rsatishda xato", show_alert=True)

async def cb_item_fav(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    uid = q.from_user.id
    _, item_id, cat_id, page = q.data.split(":")
    item_id, cat_id, page = int(item_id), int(cat_id), int(page)
    toggle_favorite(uid, item_id)
    row = get_item(uid, item_id)
    is_fav = row[7] if row else 0
    await q.answer("⭐ Sevimliga qo'shildi!" if is_fav else "💛 Sevimlilardan olib tashlandi", show_alert=True)
    update.callback_query.data = f"item_view:{item_id}:{cat_id}:{page}"
    await cb_item_view(update, ctx)

async def cb_item_del(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    uid = q.from_user.id
    _, item_id, cat_id, page = q.data.split(":")
    item_id, cat_id, page = int(item_id), int(cat_id), int(page)
    delete_item(uid, item_id)
    await q.answer("🗑 O'chirildi", show_alert=True)
    items, total = get_items(uid, cat_id, page=page)
    real_page = min(page, max(0, (total - 1) // ITEMS_PER_PAGE))
    items, total = get_items(uid, cat_id, page=real_page)
    cat = get_category(uid, cat_id)
    if cat:
        _, name, emoji = cat
        await q.edit_message_text(
            f"{emoji} <b>{name}</b>\n📦 Jami: {total} ta material",
            reply_markup=items_list_kb(uid, cat_id, real_page, total, items),
            parse_mode="HTML"
        )

async def cb_favorites(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    await q.answer()
    uid = q.from_user.id
    favs = get_favorites(uid)
    kb = InlineKeyboardMarkup([[back_btn("main_menu")]])
    if not favs:
        await q.edit_message_text("⭐ Hali sevimli material yo'q", reply_markup=kb)
        return
    lines = [f"⭐ <b>Sevimlilar</b> ({len(favs)} ta)\n"]
    for row in favs:
        _, msg_type, text, fname, cap, created, cat_name, cat_emoji, cat_id = row
        lines.append(f"• {item_preview(row)}\n   └ {cat_emoji} {cat_name}\n")
    await q.edit_message_text("\n".join(lines), reply_markup=kb, parse_mode="HTML")

async def cb_stats(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    await q.answer()
    uid = q.from_user.id
    total_cats, total_items, favorites, top_cats = get_stats(uid)
    lines = [
        "📊 <b>Statistika</b>\n",
        f"📁 Kategoriyalar: <b>{total_cats}</b>",
        f"📦 Jami materiallar: <b>{total_items}</b>",
        f"⭐ Sevimlilar: <b>{favorites}</b>",
    ]
    if top_cats:
        lines.append("\n🏆 <b>Top kategoriyalar:</b>")
        for i, (name, emoji, cnt) in enumerate(top_cats, 1):
            lines.append(f"  {i}. {emoji} {name} — {cnt} ta")
    kb = InlineKeyboardMarkup([[back_btn("main_menu")]])
    await q.edit_message_text("\n".join(lines), reply_markup=kb, parse_mode="HTML")

# ─────────────────── 🗜 COMPRESS HANDLER ─────────────────────────

async def cb_compress(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    await q.answer("⏳ Siqilmoqda…")
    uid = q.from_user.id
    parts = q.data.split(":")
    item_id, cat_id, page, level = int(parts[1]), int(parts[2]), int(parts[3]), parts[4]

    row = get_item(uid, item_id)
    if not row:
        await q.answer("❌ Topilmadi", show_alert=True)
        return
    _, msg_type, text, file_id, file_name, caption, tags, is_fav = row

    level_label = {"yuqori": "Engil (~30%)", "urtacha": "O'rtacha (~55%)", "kuchli": "Kuchli (~75%)"}
    status_msg = await q.message.reply_text(
        f"⏳ Fayl yuklanmoqda…\n🗜 Siqish darajasi: {level_label[level]}"
    )
    try:
        original_data = await download_file(ctx.bot, file_id)
        original_size = len(original_data)

        if msg_type == "photo":
            quality = QUALITY_LEVELS[level]
            compressed = await compress_image_bytes(original_data, quality)
            compressed_size = len(compressed)
            ratio = 100 - (compressed_size / original_size * 100)
            await status_msg.delete()
            await q.message.reply_photo(
                io.BytesIO(compressed),
                caption=(
                    f"🗜 <b>Siqilgan rasm</b> [{level_label[level]}]\n"
                    f"📏 Avval: <b>{human_size(original_size)}</b>\n"
                    f"📉 Keyin: <b>{human_size(compressed_size)}</b> "
                    f"(<b>-{ratio:.0f}%</b> tejaldi)"
                ),
                parse_mode="HTML"
            )

        elif msg_type in ("document", "audio", "voice", "video"):
            fname = file_name or f"fayl_{item_id}.bin"
            zip_data, zip_name = await compress_generic_bytes(original_data, fname)
            zip_size = len(zip_data)
            ratio = 100 - (zip_size / original_size * 100) if original_size else 0
            await status_msg.delete()
            await q.message.reply_document(
                io.BytesIO(zip_data),
                filename=zip_name,
                caption=(
                    f"🗜 <b>Siqilgan fayl (ZIP)</b> [{level_label[level]}]\n"
                    f"📏 Avval: <b>{human_size(original_size)}</b>\n"
                    f"📉 Keyin: <b>{human_size(zip_size)}</b> "
                    f"(<b>-{ratio:.0f}%</b> tejaldi)"
                ),
                parse_mode="HTML"
            )
        else:
            await status_msg.edit_text("❌ Bu turdagi fayl siqishni qo'llab-quvvatlamaydi")

    except Exception as e:
        logger.error(f"compress error: {e}")
        await status_msg.edit_text(f"❌ Siqishda xato: {e}")

# ─────────────────── 📦 ZIP CATEGORY HANDLER ─────────────────────

async def cb_zip_cat(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    await q.answer("⏳ ZIP yaratilmoqda…")
    uid = q.from_user.id
    cat_id = int(q.data.split(":")[1])

    cat = get_category(uid, cat_id)
    if not cat:
        await q.answer("❌ Kategoriya topilmadi", show_alert=True)
        return
    _, cat_name, cat_emoji = cat

    file_items = get_all_file_items(uid, cat_id)
    if not file_items:
        await q.answer("📭 Bu kategoriyada fayl yo'q!", show_alert=True)
        return

    status_msg = await q.message.reply_text(
        f"⏳ <b>{cat_emoji} {cat_name}</b> uchun ZIP yaratilmoqda…\n"
        f"📎 {len(file_items)} ta fayl yuklanmoqda, biroz kuting…",
        parse_mode="HTML"
    )
    try:
        zip_buf = await build_zip_for_category(ctx.bot, uid, cat_id)
        if zip_buf is None:
            await status_msg.edit_text("📭 Fayllar topilmadi")
            return
        zip_data = zip_buf.read()
        zip_size = len(zip_data)
        safe_name = "".join(c if c.isalnum() or c in (" ", "_", "-") else "_" for c in cat_name)
        zip_filename = f"{safe_name}_{datetime.now().strftime('%Y%m%d')}.zip"
        await status_msg.delete()
        await q.message.reply_document(
            io.BytesIO(zip_data),
            filename=zip_filename,
            caption=(
                f"📦 <b>{cat_emoji} {cat_name}</b> — ZIP arxiv\n"
                f"📎 {len(file_items)} ta fayl\n"
                f"📏 Hajm: <b>{human_size(zip_size)}</b>\n"
                f"📅 {datetime.now().strftime('%d.%m.%Y %H:%M')}"
            ),
            parse_mode="HTML"
        )
    except Exception as e:
        logger.error(f"zip_cat error: {e}")
        await status_msg.edit_text(f"❌ ZIP yaratishda xato: {e}")

# ─────────────────── 📦 ZIP MODE (instant) ─────────────────────

async def cb_zip_mode(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    """Foydalanuvchini ZIP rejimga o'tkazadi — yuborilgan fayllarni darhol ZIP qiladi."""
    q = update.callback_query
    await q.answer()
    ctx.user_data['state'] = 'zip_mode'
    ctx.user_data['zip_files'] = []  # to'plangan fayllar ro'yxati
    kb = InlineKeyboardMarkup([
        [InlineKeyboardButton("✅ ZIP qilib yuborish", callback_data="zip_finish")],
        [InlineKeyboardButton("❌ Bekor qilish",       callback_data="zip_cancel")],
    ])
    await q.edit_message_text(
        "📦 <b>ZIP qilish rejimi</b>\n\n"
        "Endi menga fayllarni yuboring (rasm, hujjat, video, audio).\n"
        "Bir nechta fayl yuborsangiz — hammasini bitta ZIP ga qo'shaman.\n\n"
        "Tayyor bo'lgach <b>✅ ZIP qilib yuborish</b> tugmasini bosing.",
        reply_markup=kb, parse_mode="HTML"
    )

async def cb_zip_finish(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    """To'plangan fayllarni ZIP qilib yuboradi."""
    q = update.callback_query
    await q.answer("⏳ ZIP yaratilmoqda…")
    files = ctx.user_data.get('zip_files', [])
    ctx.user_data.clear()

    if not files:
        await q.edit_message_text(
            "📭 Hech qanday fayl yuborilmadi!",
            reply_markup=main_menu_kb()
        )
        return

    status_msg = await q.message.reply_text(
        f"⏳ {len(files)} ta fayl ZIP qilinmoqda, biroz kuting…"
    )
    try:
        zip_buf = io.BytesIO()
        counters = {}
        total_original = 0
        with zipfile.ZipFile(zip_buf, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=9) as zf:
            for idx, finfo in enumerate(files, 1):
                try:
                    data = await download_file(ctx.bot, finfo['file_id'])
                except Exception as e:
                    logger.warning(f"ZIP skip: {e}")
                    continue
                total_original += len(data)
                # Arxiv ichidagi nomni aniqlash
                if finfo.get('file_name'):
                    arc_name = finfo['file_name']
                else:
                    ext = {"photo": "jpg", "video": "mp4", "audio": "mp3", "voice": "ogg"}.get(finfo['msg_type'], "bin")
                    arc_name = f"{finfo['msg_type']}_{idx}.{ext}"
                # Takrorlanishni oldini olish
                if arc_name in counters:
                    counters[arc_name] += 1
                    base, _, ext2 = arc_name.rpartition(".")
                    arc_name = f"{base}_{counters[arc_name]}.{ext2}" if ext2 else f"{arc_name}_{counters[arc_name]}"
                else:
                    counters[arc_name] = 0
                zf.writestr(arc_name, data)

        zip_buf.seek(0)
        zip_data = zip_buf.read()
        zip_size = len(zip_data)
        ratio = 100 - (zip_size / total_original * 100) if total_original else 0
        zip_filename = f"fayllar_{datetime.now().strftime('%Y%m%d_%H%M')}.zip"

        await status_msg.delete()
        await q.message.reply_document(
            io.BytesIO(zip_data),
            filename=zip_filename,
            caption=(
                f"📦 <b>ZIP arxiv tayyor!</b>\n"
                f"📎 {len(files)} ta fayl\n"
                f"📏 Avval: <b>{human_size(total_original)}</b>\n"
                f"📉 Keyin: <b>{human_size(zip_size)}</b> "
                f"(<b>-{ratio:.0f}%</b>)"
            ),
            parse_mode="HTML",
            reply_markup=InlineKeyboardMarkup([[
                InlineKeyboardButton("🏠 Menyu", callback_data="main_menu")
            ]])
        )
    except Exception as e:
        logger.error(f"zip_finish error: {e}")
        await status_msg.edit_text(f"❌ ZIP yaratishda xato: {e}")

async def cb_zip_cancel(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    await q.answer("❌ Bekor qilindi")
    ctx.user_data.clear()
    await q.edit_message_text("🏠 <b>Asosiy menyu</b>", reply_markup=main_menu_kb(), parse_mode="HTML")

# ─────────────────── SAVE / STATE HANDLERS ───────────────────────

async def cb_save_to(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    await q.answer()
    uid = q.from_user.id
    cat_id = int(q.data.split(":")[1])
    pending = ctx.user_data.pop("pending_msg", None)
    ctx.user_data.clear()
    if not pending:
        await q.edit_message_text("❌ Saqlash uchun xabar topilmadi", reply_markup=main_menu_kb())
        return
    save_item(uid, cat_id, **pending)
    cat = get_category(uid, cat_id)
    _, name, emoji = cat if cat else (None, "?", "📁")
    await q.edit_message_text(
        f"✅ <b>{emoji} {name}</b> ga saqlandi!",
        reply_markup=InlineKeyboardMarkup([[
            InlineKeyboardButton("📂 Kategoriyani ochish", callback_data=f"cat_open:{cat_id}:0"),
            InlineKeyboardButton("🏠 Menyu",               callback_data="main_menu"),
        ]]),
        parse_mode="HTML"
    )

async def cb_save_cancel(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    await q.answer("❌ Bekor qilindi")
    ctx.user_data.clear()
    await q.edit_message_text("🏠 <b>Asosiy menyu</b>", reply_markup=main_menu_kb(), parse_mode="HTML")

async def cb_noop(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    await update.callback_query.answer()

async def cb_check_sub(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    """Obuna tekshirish tugmasi."""
    q = update.callback_query
    uid = q.from_user.id
    if await is_subscribed(ctx.bot, uid):
        await q.answer("✅ Obuna tasdiqlandi!", show_alert=True)
        await q.edit_message_text(
            "🏠 <b>Asosiy menyu</b>",
            reply_markup=main_menu_kb(),
            parse_mode="HTML"
        )
    else:
        await q.answer("❌ Hali obuna bo'lmagansiz!", show_alert=True)

# ── Message handler ───────────────────────────────────────────────

async def handle_message(update: Update, ctx: ContextTypes.DEFAULT_TYPE):
    # Foydalanuvchini ro'yxatga olamiz (har xabarda last_seen yangilanadi)
    u = update.effective_user
    track_user(u.id, u.username, u.first_name)

    if not await require_subscription(update, ctx):
        return
    uid   = update.effective_user.id
    msg   = update.message
    state = ctx.user_data.get("state")

    # ── Admin broadcast rejimi ─────────────────────────────────────
    if state == "broadcast" and uid == ADMIN_ID:
        ctx.user_data.clear()
        text = msg.text or msg.caption or ""
        if not text:
            await msg.reply_text("❌ Matn yuboring (rasm/video bilan caption ham bo'ladi)")
            return

        # Barcha foydalanuvchilarni olamiz
        with get_conn() as conn:
            all_users = [r[0] for r in conn.execute("SELECT user_id FROM users").fetchall()]

        await msg.reply_text(f"⏳ {len(all_users)} ta foydalanuvchiga yuborilmoqda…")

        sent, failed = 0, 0
        for user_id in all_users:
            try:
                await ctx.bot.send_message(user_id, text, parse_mode="HTML")
                sent += 1
            except Exception:
                failed += 1
        await msg.reply_text(
            f"✅ <b>Tugadi!</b>\n"
            f"📤 Yuborildi: <b>{sent}</b>\n"
            f"❌ Xato: <b>{failed}</b>",
            parse_mode="HTML", reply_markup=admin_menu_kb() if uid == ADMIN_ID else main_menu_kb()
        )
        return

    # ── AI Prezentatsiya: mavzu kutish ─────────────────────────────
    if state == "ai_preso_topic":
        topic = (msg.text or "").strip()
        if not topic:
            await msg.reply_text("❌ Iltimos, prezentatsiya mavzusini matn ko'rinishida yozing.")
            return
        if len(topic) > 200:
            await msg.reply_text("❌ Mavzu juda uzun. Maksimum 200 belgi.")
            return
        await handle_ai_preso(update, ctx, topic)
        return

    # ── ZIP rejim: yuborilgan fayllarni to'plash ──────────────────
    if state == "zip_mode":
        m = msg
        msg_type = file_id = file_name = None

        if m.photo:
            msg_type = "photo";    file_id = m.photo[-1].file_id
        elif m.document:
            msg_type = "document"; file_id = m.document.file_id; file_name = m.document.file_name
        elif m.video:
            msg_type = "video";    file_id = m.video.file_id;    file_name = m.video.file_name
        elif m.audio:
            msg_type = "audio";    file_id = m.audio.file_id;    file_name = m.audio.file_name
        elif m.voice:
            msg_type = "voice";    file_id = m.voice.file_id
        else:
            await m.reply_text("❌ Faqat fayl yuboring (rasm, hujjat, video, audio)")
            return

        files = ctx.user_data.setdefault('zip_files', [])
        files.append({'msg_type': msg_type, 'file_id': file_id, 'file_name': file_name})

        kb = InlineKeyboardMarkup([
            [InlineKeyboardButton(f"✅ ZIP qilib yuborish ({len(files)} ta)", callback_data="zip_finish")],
            [InlineKeyboardButton("❌ Bekor qilish", callback_data="zip_cancel")],
        ])
        await m.reply_text(
            f"➕ <b>{TYPE_NAME.get(msg_type, 'Fayl')}</b> qo'shildi\n"
            f"📎 Jami: <b>{len(files)} ta</b>\n\n"
            "Yana fayl yuboring yoki ZIP qilib oling 👇",
            reply_markup=kb, parse_mode="HTML"
        )
        return

    if state == "adding_category":
        text = msg.text.strip() if msg.text else ""
        if not text:
            await msg.reply_text("❌ Matn yuboring")
            return
        emoji, name = "📁", text
        if text and ord(text[0]) > 127:
            parts = text.split(None, 1)
            if len(parts) == 2:
                emoji, name = parts
        if len(name) > 50:
            await msg.reply_text("❌ Nom juda uzun (max 50 belgi)")
            return
        ctx.user_data.clear()
        ok = add_category(uid, name, emoji)
        if ok:
            await msg.reply_text(f"✅ <b>{emoji} {name}</b> yaratildi!", reply_markup=main_menu_kb(), parse_mode="HTML")
        else:
            await msg.reply_text("⚠️ Bu nom allaqachon mavjud!", reply_markup=main_menu_kb())
        return

    if state in ("searching_all", "searching_cat"):
        query_text = msg.text.strip() if msg.text else ""
        if not query_text:
            await msg.reply_text("❌ Matn yuboring")
            return
        if state == "searching_cat":
            cat_id = ctx.user_data.get("search_cat")
            ctx.user_data.clear()
            items, total = get_items(uid, cat_id, page=0, search=query_text)
            cat = get_category(uid, cat_id)
            _, name, emoji = cat if cat else (None, "?", "📁")
            if not items:
                await msg.reply_text(
                    f"🔍 <b>«{query_text}»</b> bo'yicha {emoji} {name} da hech narsa topilmadi",
                    reply_markup=main_menu_kb(), parse_mode="HTML"
                )
            else:
                await msg.reply_text(
                    f"🔍 <b>«{query_text}»</b> — {total} ta natija ({emoji} {name})",
                    reply_markup=items_list_kb(uid, cat_id, 0, total, items), parse_mode="HTML"
                )
        else:
            ctx.user_data.clear()
            results = search_all_items(uid, query_text)
            if not results:
                await msg.reply_text(
                    f"🔍 <b>«{query_text}»</b> bo'yicha hech narsa topilmadi",
                    reply_markup=main_menu_kb(), parse_mode="HTML"
                )
            else:
                lines = [f"🔍 <b>«{query_text}»</b> — {len(results)} ta natija:\n"]
                for row in results:
                    _, msg_type, text, fname, cap, is_fav, created, cat_name, cat_emoji, cat_id = row
                    star = "⭐ " if is_fav else ""
                    lines.append(f"{star}{item_preview(row)}\n   └ {cat_emoji} {cat_name}\n")
                await msg.reply_text("\n".join(lines), reply_markup=main_menu_kb(), parse_mode="HTML")
        return

    # ── Yangi material saqlash ─────────────────────────────────────
    m = msg
    msg_type = file_id = file_name = caption = text_content = None

    if m.text and not m.text.startswith("/"):
        msg_type, text_content = "text", m.text
    elif m.photo:
        msg_type = "photo";    file_id = m.photo[-1].file_id;   caption = m.caption
    elif m.document:
        msg_type = "document"; file_id = m.document.file_id
        file_name = m.document.file_name;                        caption = m.caption
    elif m.video:
        msg_type = "video";    file_id = m.video.file_id;        caption = m.caption
    elif m.audio:
        msg_type = "audio";    file_id = m.audio.file_id;        caption = m.caption
    elif m.voice:
        msg_type = "voice";    file_id = m.voice.file_id;        caption = m.caption
    else:
        await m.reply_text("❓ Bu turdagi xabarni saqlab bo'lmaydi")
        return

    ctx.user_data['pending_msg'] = dict(
        msg_type=msg_type, text=text_content, file_id=file_id,
        file_name=file_name, caption=caption
    )
    ctx.user_data['state'] = 'saving'

    cats = get_categories(uid)
    if not cats:
        await m.reply_text(
            "📂 Hali kategoriya yo'q!\n\nAvval kategoriya yarating:",
            reply_markup=InlineKeyboardMarkup([[
                InlineKeyboardButton("➕ Kategoriya yaratish", callback_data="cat_add")
            ]])
        )
        return

    buttons = [
        [InlineKeyboardButton(f"{emoji} {name}", callback_data=f"save_to:{cat_id}")]
        for cat_id, name, emoji in cats
    ]
    buttons.append([InlineKeyboardButton("❌ Bekor qilish", callback_data="save_cancel")])
    icon  = TYPE_ICON.get(msg_type, "📎")
    tname = TYPE_NAME.get(msg_type, "Fayl")
    await m.reply_text(
        f"{icon} <b>{tname}</b> — qaysi kategoriyaga saqlash?",
        reply_markup=InlineKeyboardMarkup(buttons), parse_mode="HTML"
    )

# ─────────────────────── WEB SERVER (Render uchun) ────────────────
# Render hosting platforma bot ishlaydigan portga so'rov yuborib turadi.
# UptimeRobot ham shu URL ga so'rov yuborib botni "uxlamasligini" ta'minlaydi.

flask_app = Flask(__name__)

@flask_app.route('/')
def home():
    return "🤖 EduSave Bot ishlamoqda!"

@flask_app.route('/health')
def health():
    return "OK"

def run_flask():
    port = int(os.environ.get('PORT', 8080))
    flask_app.run(host='0.0.0.0', port=port)

# ─────────────────────────── MAIN ────────────────────────────────

def main():
    token = os.getenv("BOT_TOKEN")
    if not token:
        print("❌ BOT_TOKEN muhit o'zgaruvchisi topilmadi!")
        print("   export BOT_TOKEN=your_token_here")
        return
    init_db()
    logger.info("DB initialized ✅")

    # Flask web serverni alohida thread da ishga tushiramiz (Render uchun)
    flask_thread = threading.Thread(target=run_flask, daemon=True)
    flask_thread.start()
    logger.info("🌐 Web server ishga tushdi (port 8080)")
    app = Application.builder().token(token).build()
    app.add_handler(CommandHandler("start", cmd_start))
    app.add_handler(CommandHandler("menu",  cmd_start))
    app.add_handler(CommandHandler("admin", cmd_admin))
    cb_map = {
        "^main_menu$":         cb_main_menu,
        "^cat_list$":          cb_cat_list,
        "^cat_add$":           cb_cat_add,
        "^cat_open:":          cb_cat_open,
        "^cat_del_confirm:":   cb_cat_del_confirm,
        "^cat_del:":           cb_cat_del,
        "^cat_search:":        cb_cat_search,
        "^search_start$":      cb_search_start,
        "^item_view:":         cb_item_view,
        "^item_fav:":          cb_item_fav,
        "^item_del:":          cb_item_del,
        "^favorites$":         cb_favorites,
        "^stats$":             cb_stats,
        "^save_to:":           cb_save_to,
        "^save_cancel$":       cb_save_cancel,
        "^compress:":          cb_compress,
        "^zip_cat:":           cb_zip_cat,
        "^zip_mode$":          cb_zip_mode,
        "^zip_finish$":        cb_zip_finish,
        "^zip_cancel$":        cb_zip_cancel,
        "^check_sub$":         cb_check_sub,
        "^admin_users$":       cb_admin_users,
        "^admin_stats$":       cb_admin_stats,
        "^admin_recent$":      cb_admin_recent,
        "^admin_back$":        cb_admin_back,
        "^admin_broadcast$":   cb_admin_broadcast,
        "^ai_preso$":          cb_ai_preso,
        "^noop$":              cb_noop,
    }
    for pattern, handler in cb_map.items():
        # check_sub o'zi obunani tekshiradi, qolganlar uchun wrap qilamiz
        if pattern == "^check_sub$":
            app.add_handler(CallbackQueryHandler(handler, pattern=pattern))
        else:
            # Har bir callback'dan oldin obunani tekshiramiz
            def make_protected(h):
                async def protected(update, ctx):
                    if not await require_subscription(update, ctx):
                        return
                    return await h(update, ctx)
                return protected
            app.add_handler(CallbackQueryHandler(make_protected(handler), pattern=pattern))
    app.add_handler(MessageHandler(
        filters.TEXT | filters.PHOTO | filters.Document.ALL |
        filters.VIDEO | filters.AUDIO | filters.VOICE,
        handle_message
    ))
    logger.info("🤖 Bot ishga tushdi!")
    app.run_polling(allowed_updates=Update.ALL_TYPES)

if __name__ == "__main__":
    import asyncio
    # Python 3.14 da event loop ni yangidan yaratamiz
    try:
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
    except Exception:
        pass
    main()
